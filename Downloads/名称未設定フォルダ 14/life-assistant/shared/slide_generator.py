"""
スライド生成モジュール
議事録からINTLOOPテンプレート準拠のPowerPointスライドを生成
"""
import json
import io
import os
from typing import Dict, Any, List, Optional
from io import BytesIO
from pptx import Presentation
from pptx.util import Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from .config import Config
from .claude_client import ClaudeClient


# テンプレートカラーパレット
C_PRIMARY = RGBColor(0x00, 0x39, 0x65)  # ロゴ基調
C_SECONDARY = RGBColor(0x00, 0x86, 0xC5)  # 強調ブルー
C_ACCENT = RGBColor(0xED, 0x51, 0x4E)  # ポイント赤
C_BLACK = RGBColor(0x00, 0x00, 0x00)
C_GRAY = RGBColor(0x59, 0x59, 0x59)
C_LIGHT_GRAY = RGBColor(0xD9, 0xD9, 0xD9)
C_BG_LIGHT = RGBColor(0xF7, 0xF7, 0xF7)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_BODY = "BIZ UDPゴシック"
FONT_HEADING = "BIZ UDPゴシック"
FONT_NUMBER = "Helvetica"

# レイアウト名（テンプレートに存在するもの）
LAYOUT_NAMES = {
    "title": "タイトルスライド_クライアント向け",
    "toc": "目次スライド",
    "section": "サブタイトルページ",
    "basic": "基本スライド",
    "basic_no_line": "基本：タイトル下線なし",
    "table_insert": "基本：表テーブル挿入スライド",
    "box3_1": "基本：3BOX①",
    "box4": "基本：4BOX",
    "flow": "基本：フロー図",
    "logo": "ロゴページ",
}


def _int_emu(*values):
    """EMU値を整数に変換"""
    return tuple(int(v) for v in values)


class Decorator:
    """スライドに装飾要素を追加するヘルパー"""

    @staticmethod
    def add_rounded_rect(slide, left, top, width, height,
                         fill_color=None, line_color=None, line_width=Pt(1.5),
                         text="", font_size=11, font_color=C_WHITE,
                         font_bold=True, align=PP_ALIGN.CENTER,
                         vertical=MSO_ANCHOR.MIDDLE):
        """角丸四角形"""
        left, top, width, height = _int_emu(left, top, width, height)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = 0.08

        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        else:
            shape.fill.background()

        ln = shape.line
        if line_color:
            ln.fill.solid()
            ln.fill.fore_color.rgb = line_color
            ln.width = line_width
        else:
            ln.fill.background()

        if text:
            tf = shape.text_frame
            tf.word_wrap = True
            tf.margin_left = Cm(0.3)
            tf.margin_right = Cm(0.3)
            tf.margin_top = Cm(0.2)
            tf.margin_bottom = Cm(0.2)
            p = tf.paragraphs[0]
            p.alignment = align
            p.space_after = Pt(0)
            run = p.add_run()
            run.text = text
            run.font.name = FONT_BODY
            run.font.size = Pt(font_size)
            run.font.color.rgb = font_color
            run.font.bold = font_bold

        return shape

    @staticmethod
    def add_number_badge(slide, number, left, top,
                          size=Cm(1.2), bg_color=C_PRIMARY):
        """番号バッジ（丸の中に数字）"""
        left, top, size = _int_emu(left, top, size)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, left, top, size, size)
        circle.fill.solid()
        circle.fill.fore_color.rgb = bg_color
        circle.line.fill.background()
        tf = circle.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(number)
        run.font.name = FONT_NUMBER
        run.font.size = Pt(11)
        run.font.color.rgb = C_WHITE
        run.font.bold = True
        bodyPr = circle._element.txBody.find(qn('a:bodyPr'))
        bodyPr.set('anchor', 'ctr')

    @staticmethod
    def add_type_badge(slide, text, left=Cm(29.5), top=Cm(0.6),
                       width=Cm(3.5), height=Cm(0.8),
                       bg_color=C_SECONDARY, font_color=C_WHITE):
        """セクションタイプのカラーバッジ（右上）"""
        left, top, width, height = _int_emu(left, top, width, height)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = 0.25
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = False
        tf.margin_left = Cm(0.1)
        tf.margin_right = Cm(0.1)
        tf.margin_top = Cm(0)
        tf.margin_bottom = Cm(0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = FONT_BODY
        run.font.size = Pt(9)
        run.font.color.rgb = font_color
        run.font.bold = True
        bodyPr = shape._element.txBody.find(qn('a:bodyPr'))
        bodyPr.set('anchor', 'ctr')
        return shape


class SlideGenerator:
    """議事録からスライドを生成（INTLOOPテンプレート準拠）"""

    TEMPLATE_PATH = os.path.join(
        os.path.dirname(os.path.abspath(__file__)),
        "templates",
        "INTLOOP_スライドテンプレート_コピー.pptx"
    )

    def __init__(self, config: Optional[Config] = None):
        if config is None:
            from .config import get_config
            config = get_config()
        self.config = config
        self.claude = ClaudeClient(config)
        self.deco = Decorator()
        self._template_exists = os.path.exists(self.TEMPLATE_PATH)

    def _init_presentation(self):
        """プレゼンテーションを初期化"""
        if self._template_exists:
            prs = Presentation(self.TEMPLATE_PATH)
            # テンプレートの既存スライドを削除
            while len(prs.slides) > 0:
                rId = prs.slides._sldIdLst[0].rId
                prs.part.drop_rel(rId)
                prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
            return prs
        else:
            # テンプレートがない場合は空のプレゼンテーション
            prs = Presentation()
            prs.slide_width = Cm(33.867)
            prs.slide_height = Cm(19.05)
            return prs

    def _get_layout(self, prs, layout_key):
        """レイアウトを取得"""
        layout_name = LAYOUT_NAMES.get(layout_key, LAYOUT_NAMES["basic"])
        for master in prs.slide_masters:
            for layout in master.slide_layouts:
                if layout.name == layout_name:
                    return layout
        # フォールバック
        return prs.slide_masters[0].slide_layouts[0]

    def _get_ph(self, slide, idx):
        """プレースホルダーを取得"""
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == idx:
                return ph
        return None

    def _set_ph_text(self, slide, idx, text):
        """プレースホルダーにテキストを設定"""
        ph = self._get_ph(slide, idx)
        if ph and ph.has_text_frame:
            ph.text = text
        return ph

    def _set_ph_lines(self, slide, idx, lines):
        """プレースホルダーに複数行テキストを設定"""
        ph = self._get_ph(slide, idx)
        if not ph or not ph.has_text_frame:
            return ph
        tf = ph.text_frame
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
                p.text = line
            else:
                p = tf.add_paragraph()
                p.text = line
            p.space_after = Pt(6)
        return ph

    def parse_meeting_notes(self, text: str) -> Dict[str, Any]:
        """議事録テキストを解析してスライド用の構造化データに変換"""
        prompt = f"""以下の議事録を分析し、スライド作成用の構造化データに変換してください。

## 議事録
{text}

## 出力形式（JSON）
{{
  "title": "会議タイトル（改行可）",
  "subtitle": "サブタイトル・プロジェクト名等",
  "date": "YYYY年M月D日",
  "sections": [
    {{
      "heading": "セクション見出し",
      "subhead": "サブ見出し（任意）",
      "type": "セクション種別",
      "key_points": ["要点1", "要点2", ...],
      "action_items": [
        {{"assignee": "担当者", "action": "内容", "deadline": "期限"}}
      ]
    }}
  ]
}}

## セクション種別（type）
- overview: プロジェクト概要、背景、目的
- team_structure: 体制、メンバー、役割分担
- technology: 技術方針、システム構成
- schedule: スケジュール、マイルストーン
- risk: リスク、課題、懸念事項
- action_items: アクションアイテム、TODO
- budget: 予算、コスト
- discussion: 議論、検討事項
- general: 上記に当てはまらないもの

注意：
- 要点は各セクション3〜7個に絞り、簡潔に
- 決定事項は【決定】を先頭に付ける
- 重要な情報は必ず含める

JSONのみを出力してください。"""

        response = self.claude.client.messages.create(
            model=self.claude.model,
            max_tokens=4096,
            messages=[{"role": "user", "content": prompt}]
        )

        content = response.content[0].text
        result = self.claude._parse_json_response(content)

        if not result:
            return {
                "title": "会議議事録",
                "subtitle": "",
                "date": "",
                "sections": [{
                    "heading": "議論内容",
                    "type": "general",
                    "key_points": [text[:200]]
                }]
            }

        return result

    def generate_pptx(self, structured_data: Dict[str, Any]) -> bytes:
        """構造化データからPowerPointファイルを生成"""
        prs = self._init_presentation()

        # タイトルスライド
        self._add_title_slide(prs, structured_data)

        sections = structured_data.get("sections", [])

        # 目次スライド
        if len(sections) > 1:
            self._add_toc_slide(prs, sections)

        # 各セクション
        total = len(sections)
        for si, section in enumerate(sections, 1):
            heading = section.get("heading", "")
            subhead = section.get("subhead", "")
            section_type = section.get("type", "general")
            key_points = section.get("key_points", [])
            action_items = section.get("action_items", [])

            # セクション区切りスライド
            if total > 1:
                self._add_section_slide(prs, heading, section_num=si, total_sections=total)

            # コンテンツスライド
            if action_items:
                self._add_action_items_slide(prs, heading, action_items)
            elif key_points:
                self._add_content_slide(prs, heading, subhead, key_points, section_type)

        # ロゴページ
        self._add_logo_slide(prs)

        # バイトデータとして出力
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()

    def _add_title_slide(self, prs, data):
        """タイトルスライドを追加"""
        if self._template_exists:
            layout = self._get_layout(prs, "title")
            slide = prs.slides.add_slide(layout)
            self._set_ph_text(slide, 0, data.get("title", "会議議事録"))
            self._set_ph_text(slide, 10, data.get("subtitle", ""))
            self._set_ph_text(slide, 13, data.get("date", ""))
        else:
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            # タイトル
            title_box = slide.shapes.add_textbox(Cm(1), Cm(6), Cm(31), Cm(3))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = data.get("title", "会議議事録")
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = C_PRIMARY
            p.alignment = PP_ALIGN.CENTER
            # サブタイトル
            if data.get("subtitle"):
                sub_box = slide.shapes.add_textbox(Cm(1), Cm(10), Cm(31), Cm(1.5))
                tf2 = sub_box.text_frame
                p2 = tf2.paragraphs[0]
                p2.text = data["subtitle"]
                p2.font.size = Pt(24)
                p2.font.color.rgb = C_GRAY
                p2.alignment = PP_ALIGN.CENTER
            # 日付
            if data.get("date"):
                date_box = slide.shapes.add_textbox(Cm(1), Cm(12), Cm(31), Cm(1))
                tf3 = date_box.text_frame
                p3 = tf3.paragraphs[0]
                p3.text = data["date"]
                p3.font.size = Pt(18)
                p3.font.color.rgb = C_GRAY
                p3.alignment = PP_ALIGN.CENTER
        return slide

    def _add_toc_slide(self, prs, sections):
        """目次スライドを追加"""
        if self._template_exists:
            layout = self._get_layout(prs, "toc")
            slide = prs.slides.add_slide(layout)
            lines = [f"   {s['heading']}" for s in sections]
            self._set_ph_lines(slide, 13, lines)
            # 番号バッジ
            ph = self._get_ph(slide, 13)
            if ph:
                badge_size = Cm(0.8) if len(sections) > 6 else Cm(0.9)
                base_left = int(ph.left - Cm(0.1))
                base_top = int(ph.top + Cm(0.05))
                line_height = int(ph.height / max(len(sections), 1))
                for i in range(len(sections)):
                    self.deco.add_number_badge(
                        slide, i + 1,
                        base_left, base_top + line_height * i,
                        size=badge_size, bg_color=C_SECONDARY)
        else:
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            # タイトル
            title_box = slide.shapes.add_textbox(Cm(1), Cm(1), Cm(31), Cm(2))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = "目次"
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = C_PRIMARY
            # 項目
            content_box = slide.shapes.add_textbox(Cm(3), Cm(4), Cm(27), Cm(12))
            tf2 = content_box.text_frame
            for i, section in enumerate(sections):
                para = tf2.add_paragraph() if i > 0 else tf2.paragraphs[0]
                para.text = f"{i+1}. {section['heading']}"
                para.font.size = Pt(20)
                para.font.color.rgb = C_BLACK
                para.space_after = Pt(12)
        return slide

    def _add_section_slide(self, prs, heading, section_num=0, total_sections=0):
        """セクション区切りスライドを追加"""
        if self._template_exists:
            layout = self._get_layout(prs, "section")
            slide = prs.slides.add_slide(layout)
            self._set_ph_text(slide, 10, heading)
            # セクション番号
            if section_num > 0:
                txBox = slide.shapes.add_textbox(Cm(0.8), Cm(13.5), Cm(3.0), Cm(2.5))
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = f"{section_num:02d}"
                run.font.name = FONT_NUMBER
                run.font.size = Pt(48)
                run.font.color.rgb = C_LIGHT_GRAY
        else:
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            # 背景色
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = C_PRIMARY
            # タイトル
            title_box = slide.shapes.add_textbox(Cm(2), Cm(7), Cm(29), Cm(3))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = heading
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.color.rgb = C_WHITE
            p.alignment = PP_ALIGN.CENTER
            # 番号
            if section_num > 0:
                num_box = slide.shapes.add_textbox(Cm(1), Cm(15), Cm(3), Cm(2))
                tf2 = num_box.text_frame
                p2 = tf2.paragraphs[0]
                p2.text = f"{section_num:02d}"
                p2.font.size = Pt(36)
                p2.font.color.rgb = C_LIGHT_GRAY
        return slide

    def _add_content_slide(self, prs, heading, subhead, key_points, section_type="general"):
        """コンテンツスライドを追加"""
        if self._template_exists:
            layout = self._get_layout(prs, "table_insert")
            slide = prs.slides.add_slide(layout)
            self._set_ph_text(slide, 0, heading)
            if subhead:
                self._set_ph_text(slide, 11, subhead)
            # 本文
            bullet_lines = []
            for point in key_points:
                if point.startswith("【"):
                    bullet_lines.append(point)
                else:
                    bullet_lines.append(f"• {point}")
            ph = self._get_ph(slide, 12)
            if ph and ph.has_text_frame:
                tf = ph.text_frame
                for i, line in enumerate(bullet_lines):
                    if i == 0:
                        p = tf.paragraphs[0]
                        p.text = line
                    else:
                        p = tf.add_paragraph()
                        p.text = line
                    p.space_after = Pt(8)
                    p.line_spacing = 1.2
            # セクションタイプバッジ
            type_label = self._section_type_label(section_type)
            if type_label:
                self.deco.add_type_badge(slide, type_label, left=Cm(29.5), top=Cm(0.6))
        else:
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            # タイトル
            title_box = slide.shapes.add_textbox(Cm(1), Cm(1), Cm(31), Cm(1.5))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = heading
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = C_PRIMARY
            # サブヘッド
            if subhead:
                sub_box = slide.shapes.add_textbox(Cm(1), Cm(2.8), Cm(31), Cm(1))
                tf2 = sub_box.text_frame
                p2 = tf2.paragraphs[0]
                p2.text = subhead
                p2.font.size = Pt(16)
                p2.font.color.rgb = C_GRAY
            # 本文
            content_box = slide.shapes.add_textbox(Cm(1.5), Cm(4.5), Cm(30), Cm(13))
            tf3 = content_box.text_frame
            tf3.word_wrap = True
            for i, point in enumerate(key_points):
                para = tf3.add_paragraph() if i > 0 else tf3.paragraphs[0]
                if point.startswith("【"):
                    para.text = point
                else:
                    para.text = f"• {point}"
                para.font.size = Pt(18)
                para.font.color.rgb = C_BLACK
                para.space_after = Pt(10)
        return slide

    def _add_action_items_slide(self, prs, heading, items):
        """アクションアイテムスライドを追加"""
        if self._template_exists:
            layout = self._get_layout(prs, "table_insert")
            slide = prs.slides.add_slide(layout)
            self._set_ph_text(slide, 0, heading)
            # テーブル作成
            rows = len(items) + 1
            cols = 4
            table_shape = slide.shapes.add_table(
                rows, cols, Cm(1.5), Cm(4.5), Cm(29.0), Cm(1.2 + 1.3 * len(items)))
            table = table_shape.table
            table.columns[0].width = Cm(1.8)
            table.columns[1].width = Cm(4.5)
            table.columns[2].width = Cm(17.5)
            table.columns[3].width = Cm(5.2)
            # ヘッダー
            headers = ["No.", "担当", "アクション", "期限"]
            for ci, header in enumerate(headers):
                cell = table.cell(0, ci)
                cell.text = header
                for para in cell.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.CENTER
                    for run in para.runs:
                        run.font.name = FONT_HEADING
                        run.font.size = Pt(11)
                        run.font.bold = True
                        run.font.color.rgb = C_WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_PRIMARY
            # データ行
            for ri, item in enumerate(items, 1):
                vals = [str(ri),
                        item.get("assignee", ""),
                        item.get("action", ""),
                        item.get("deadline", "")]
                for ci, val in enumerate(vals):
                    cell = table.cell(ri, ci)
                    cell.text = val
                    for para in cell.text_frame.paragraphs:
                        para.alignment = PP_ALIGN.CENTER if ci != 2 else PP_ALIGN.LEFT
                        for run in para.runs:
                            run.font.name = FONT_BODY
                            run.font.size = Pt(11)
                            run.font.color.rgb = C_BLACK
                # ストライプ
                if ri % 2 == 0:
                    for ci in range(cols):
                        table.cell(ri, ci).fill.solid()
                        table.cell(ri, ci).fill.fore_color.rgb = C_BG_LIGHT
                else:
                    for ci in range(cols):
                        table.cell(ri, ci).fill.solid()
                        table.cell(ri, ci).fill.fore_color.rgb = C_WHITE
            # バッジ
            self.deco.add_type_badge(slide, "アクション", left=Cm(29.5), top=Cm(0.6))
        else:
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            # タイトル
            title_box = slide.shapes.add_textbox(Cm(1), Cm(1), Cm(31), Cm(1.5))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = heading
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = C_PRIMARY
            # 内容
            content_box = slide.shapes.add_textbox(Cm(1.5), Cm(3.5), Cm(30), Cm(14))
            tf2 = content_box.text_frame
            tf2.word_wrap = True
            for i, item in enumerate(items):
                para = tf2.add_paragraph() if i > 0 else tf2.paragraphs[0]
                task = item.get("action", "")
                assignee = item.get("assignee", "")
                deadline = item.get("deadline", "")
                parts = [f"• {task}"]
                if assignee:
                    parts.append(f"【担当: {assignee}】")
                if deadline:
                    parts.append(f"（期限: {deadline}）")
                para.text = " ".join(parts)
                para.font.size = Pt(16)
                para.font.color.rgb = C_BLACK
                para.space_after = Pt(12)
        return slide

    def _add_logo_slide(self, prs):
        """ロゴページを追加"""
        if self._template_exists:
            layout = self._get_layout(prs, "logo")
            slide = prs.slides.add_slide(layout)
        else:
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            # 中央にテキスト
            text_box = slide.shapes.add_textbox(Cm(1), Cm(8), Cm(31), Cm(3))
            tf = text_box.text_frame
            p = tf.paragraphs[0]
            p.text = "Thank You"
            p.font.size = Pt(48)
            p.font.color.rgb = C_PRIMARY
            p.alignment = PP_ALIGN.CENTER
        return slide

    def _section_type_label(self, section_type):
        """セクション種別 → 日本語バッジラベル"""
        labels = {
            "overview": "概要",
            "team_structure": "体制",
            "technology": "技術",
            "schedule": "スケジュール",
            "risk": "リスク・課題",
            "action_items": "アクション",
            "budget": "予算",
            "discussion": "議論",
            "general": "情報",
        }
        return labels.get(section_type, "")

    def generate_from_text(self, text: str) -> bytes:
        """議事録テキストからPPTXを生成する統合メソッド"""
        meeting_data = self.parse_meeting_notes(text)
        return self.generate_pptx(meeting_data)
